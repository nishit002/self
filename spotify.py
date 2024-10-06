import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Function to create the PowerPoint file
def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    title_1 = slide_1.shapes.title
    subtitle_1 = slide_1.placeholders[1]
    title_1.text = "Spotify: The Challenges of an Online Music Service"
    subtitle_1.text = "Legal and Profitable\nPresented by: Your Name\nDate"

    # Slide 2: Introduction
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    title_2 = slide_2.shapes.title
    title_2.text = "Introduction"
    
    text_box = slide_2.shapes[1].text_frame
    text_box.clear()
    text_box.text = "Spotify was founded in 2006 by Daniel Ek and Martin Lorentzon in Sweden."

    p = text_box.add_paragraph()
    p.text = "• Aimed to address illegal music downloads."
    p = text_box.add_paragraph()
    p.text = "• Offers over 30 million songs for streaming."
    p = text_box.add_paragraph()
    p.text = "• Key challenges include securing licensing agreements and achieving profitability."
    p = text_box.add_paragraph()
    p.text = "• Despite offering both free and premium services, profitability has remained a challenge due to high licensing fees."

    # Slide 3: Key Financials (2013-2014)
    slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
    title_3 = slide_3.shapes.title
    title_3.text = "Key Financials (2013-2014)"
    
    # Add a table for financial data
    rows = 5
    cols = 3
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9.0)
    height = Inches(2.0)

    table = slide_3.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(3.5)

    # Table content
    table.cell(0, 0).text = "Financial Metric"
    table.cell(0, 1).text = "2013"
    table.cell(0, 2).text = "2014"
    
    table.cell(1, 0).text = "Revenue (€ millions)"
    table.cell(1, 1).text = "747"
    table.cell(1, 2).text = "1,080"
    
    table.cell(2, 0).text = "Net Loss (€ millions)"
    table.cell(2, 1).text = "93"
    table.cell(2, 2).text = "162"
    
    table.cell(3, 0).text = "Total Accumulated Losses (€ millions)"
    table.cell(3, 1).text = "200"
    table.cell(3, 2).text = "262"

    table.cell(4, 0).text = "Subscription Revenue Contribution"
    table.cell(4, 1).text = "91%"
    table.cell(4, 2).text = "Still dominated by Premium (~90%)"

    # Slide 4: Monetization Strategy and Freemium Model
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    title_4 = slide_4.shapes.title
    title_4.text = "Monetization Strategy and Freemium Model"
    
    text_frame_4 = slide_4.shapes[1].text_frame
    text_frame_4.clear()
    text_frame_4.text = "Spotify operates under a freemium model."
    
    p = text_frame_4.add_paragraph()
    p.text = "• Free with ads or €9.99/month for Premium without ads."
    p = text_frame_4.add_paragraph()
    p.text = "• Premium accounts for 91% of total revenue."
    p = text_frame_4.add_paragraph()
    p.text = "• Only 20-27% of users convert to Premium."
    p = text_frame_4.add_paragraph()
    p.text = "• Heavy reliance on Premium for profitability, as advertising revenue remains low."

    # Slide 5: Competitors
    slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
    title_5 = slide_5.shapes.title
    title_5.text = "Spotify's Competitors"
    
    text_frame_5 = slide_5.shapes[1].text_frame
    text_frame_5.clear()
    text_frame_5.text = "Spotify faces competition from major streaming services, including:"
    
    p = text_frame_5.add_paragraph()
    p.text = "• Apple Music (€9.99/month), Pandora (€4.99/month), Deezer (€9.99/month), Tidal (€19.99/month)."
    p = text_frame_5.add_paragraph()
    p.text = "• Spotify maintains an edge through social integration (Facebook, Twitter) and personalized features like Discover Weekly."

    # Slide 6: Financial Challenges and Path to Profitability
    slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
    title_6 = slide_6.shapes.title
    title_6.text = "Financial Challenges and Path to Profitability"
    
    text_frame_6 = slide_6.shapes[1].text_frame
    text_frame_6.clear()
    text_frame_6.text = "Challenges include:"
    
    p = text_frame_6.add_paragraph()
    p.text = "• High licensing fees, accounting for 70% of revenue."
    p = text_frame_6.add_paragraph()
    p.text = "• Spotify pays out 70% of its revenue to music rights holders, meaning costs grow with revenue."
    p = text_frame_6.add_paragraph()
    p.text = "• Growing competition from Apple Music, which offers exclusive artist deals."
    p = text_frame_6.add_paragraph()
    p.text = "• Path to profitability requires converting more free users to Premium subscriptions."
    p = text_frame_6.add_paragraph()
    p.text = "• Negotiating better licensing deals with record labels and increasing ad revenue from free users."

    # Slide 7: Conclusion
    slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
    title_7 = slide_7.shapes.title
    title_7.text = "Conclusion"
    
    text_frame_7 = slide_7.shapes[1].text_frame
    text_frame_7.clear()
    text_frame_7.text = "Spotify has built a highly successful platform, but faces significant challenges in achieving profitability."
    
    p = text_frame_7.add_paragraph()
    p.text = "• Licensing fees and revenue-sharing agreements weigh heavily on profitability."
    p = text_frame_7.add_paragraph()
    p.text = "• Spotify's success lies in continuing to innovate and convert more users to paid subscriptions."
    p = text_frame_7.add_paragraph()
    p.text = "• In the future, better licensing deals and expanding into new regions will be crucial."

    return prs

# Streamlit App Title
st.title("Spotify: The Challenges of an Online Music Service")

# Introduction Section
st.markdown("""
## Slide 2: Introduction
Spotify was founded in 2006 by Daniel Ek and Martin Lorentzon in Sweden. The service aimed to address the growing issue of illegal music downloads by offering a legal and user-friendly alternative for consumers.
Challenges include securing licensing agreements, profitability, and fair compensation for artists.
""")

# Financial Data (2013-2014)
financial_data = {
    'Year': [2013, 2014],
    'Revenue (in € millions)': [747, 1080],
    'Net Loss (in € millions)': [-93, -162],
    'Accumulated Losses (in € millions)': [200, 262]
}

df = pd.DataFrame(financial_data)

# Display Financial Data Table
st.markdown("## Slide 3: Key Financials (2013-2014)")
st.dataframe(df)

# Revenue Growth Chart
st.markdown("### Spotify Revenue Growth (2013-2014)")
fig, ax = plt.subplots()
ax.plot(df['Year'], df['Revenue (in € millions)'], marker='o', color='green')
ax.set_xlabel('Year')
ax.set_ylabel('Revenue (in € millions)')
st.pyplot(fig)

# Net Loss Chart
st.markdown("### Spotify Net Loss (2013-2014)")
fig, ax = plt.subplots()
ax.bar(df['Year'], df['Net Loss (in € millions)'], color='red')
ax.set_xlabel('Year')
ax.set_ylabel('Net Loss (in € millions)')
st.pyplot(fig)

# Error Handling for PowerPoint Download
try:
    # Download PowerPoint Presentation
    st.markdown("## Download the Detailed PowerPoint Presentation")
    prs = create_presentation()
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    st.download_button(label="Download Presentation", data=pptx_io, file_name="Spotify_Detailed_Presentation_v5.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
except Exception as e:
    st.error(f"An error occurred while generating the PowerPoint: {e}")
